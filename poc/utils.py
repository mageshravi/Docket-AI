from collections.abc import Generator

import tiktoken
from django.conf import settings
from openai import OpenAI


def create_chunks_for_vector_embedding(text_content: str) -> list:
    """
    Splits the text content into chunks of approximately 8000 tokens.

    Args:
        text_content (str): The text content to be split into chunks.

    Returns:
        list: A list of text chunks.
    """
    chunk_size = 8000

    encoding = tiktoken.encoding_for_model("gpt-4o")
    tokens = encoding.encode(text_content)

    if len(tokens) <= chunk_size:
        return [text_content]

    chunks = []
    # split into chunks of approximately 8000 tokens with some overlap
    for i in range(0, len(tokens), chunk_size):
        chunk = tokens[i : i + chunk_size]

        if i > 0:
            # Add overlap of 100 tokens
            chunk = tokens[i - 100 : i + chunk_size]

        text_chunk = encoding.decode(chunk)
        chunks.append(text_chunk)

    return chunks


def create_vector_embedding(chunks: list) -> list[dict]:
    """
    Creates vector embeddings for the provided text chunks using OpenAI's API.

    Args:
        chunks (list): A list of text chunks to be embedded.

    Raises:
        OpenAIError: If there is an error with the OpenAI API.

    Returns:
        list: A list of dictionaries containing the text and its corresponding embedding.
    """
    openai = OpenAI(api_key=settings.OPENAI_API_KEY)
    embeddings = []

    for chunk in chunks:
        response = openai.embeddings.create(input=chunk, model="text-embedding-3-small")
        embeddings.append({"text": chunk, "embedding": response.data[0].embedding})

    return embeddings


def extract_text_from_pdf(file_path: str) -> Generator[str, None, None]:
    """
    Extracts text from each page of a PDF file.

    Args:
        file_path (str): The path to the PDF file.

    Returns:
        Generator[str, None, None]: A generator that yields text from the PDF by page.
    """
    import pypdf

    reader = pypdf.PdfReader(file_path)

    for page in reader.pages:
        yield page.extract_text() or ""


def extract_text_from_xlsx(file_path: str) -> Generator[str, None, None]:
    """
    Extracts text from each sheet of an XLSX file.

    Args:
        file_path (str): The path to the XLSX file.

    Returns:
        Generator[str, None, None]: A generator that yields text from the XLSX by sheet.
    """
    import pandas as pd

    xls = pd.ExcelFile(file_path)
    for sheet_name in xls.sheet_names:
        df = pd.read_excel(xls, sheet_name)
        text_content = "Sheet: " + sheet_name + "\n"
        for _, row in df.iterrows():
            text_content += (
                "\t".join(str(value) for value in row if pd.notna(value)) + "\n"
            )

        yield text_content


def extract_text_from_docx(file_path: str) -> Generator[str, None, None]:
    """
    Extracts text from a DOCX file.

    Args:
        file_path (str): The path to the DOCX file.

    Returns:
        Generator[str, None, None]: A generator that yields text from the DOCX file by block.
    """
    from docx import Document
    from docx.oxml.table import CT_Tbl
    from docx.oxml.text.paragraph import CT_P
    from docx.table import Table

    doc = Document(file_path)
    for child in doc.element.body.iterchildren():
        if isinstance(child, CT_P):
            yield child.text
        elif isinstance(child, CT_Tbl):
            table = Table(child, doc)
            text = ""
            for row in table.rows:
                text += "\t".join(cell.text for cell in row.cells) + "\n"

            yield text


def extract_text_from_pptx(file_path: str) -> Generator[str, None, None]:
    """
    Extracts text from a PPTX file.

    Args:
        file_path (str): The path to the PPTX file.

    Returns:
        Generator[str, None, None]: A generator that yields text from the PPTX by slide.
    """
    from pptx import Presentation

    presentation = Presentation(file_path)
    for slide in presentation.slides:
        text_content = f"Slide {slide.slide_id}:\n"
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_content += shape.text + "\n"

        yield text_content.strip()


def extract_text_from_txt(file_path: str) -> Generator[str, None, None]:
    """
    Extracts text from a TXT file.

    Args:
        file_path (str): The path to the TXT file.

    Raises:
        FileNotFoundError: If the file does not exist.

    Returns:
        Generator[str, None, None]: A generator that yields text from the TXT file in chunks of 8000 characters.
    """
    BUFFER_SIZE = 8000  # Size of each chunk to yield
    buffer = ""
    with open(file_path, "r", encoding="utf-8") as file:
        for line in file:
            buffer += line
            if len(buffer) >= BUFFER_SIZE:
                yield buffer
                buffer = ""

    if buffer:
        yield buffer


def extract_text_from_csv(file_path: str) -> Generator[str, None, None]:
    """
    Extracts text from a CSV file.

    Args:
        file_path (str): The path to the CSV file.

    Returns:
        Generator[str, None, None]: A generator that yields text from the CSV by chunks of 200 rows.
    """
    import pandas as pd

    df = pd.read_csv(file_path)
    chunk = []
    for _, row in df.iterrows():
        chunk.append("\t".join(str(value) for value in row if pd.notna(value)))
        if len(chunk) >= 200:
            yield "\n".join(chunk)
            chunk = []

    if chunk:
        yield "\n".join(chunk)
